"""Build apresentacao.pptx from the project decisions/architecture."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x1A, 0x44, 0x80)
DARK = RGBColor(0x22, 0x22, 0x22)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xEE, 0xEE, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def add_text(slide, left, top, width, height, text, *, size=18, bold=False,
             color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return tb


def add_bullets(slide, left, top, width, height, items, *, size=18, color=DARK):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            level, txt = item
        else:
            level, txt = 0, item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        bullet = "• " if level == 0 else "– "
        run = p.add_run()
        run.text = bullet + txt
        run.font.size = Pt(size - level * 2)
        run.font.color.rgb = color
        run.font.name = "Calibri"
        p.space_after = Pt(4)
    return tb


def add_header(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    add_text(slide, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7),
             title, size=30, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.4),
                 subtitle, size=16, color=GREY)


def add_footer(slide, page):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(8), Inches(0.3),
             "RAG ANEEL · Desafio CEIA", size=10, color=GREY)
    add_text(slide, Inches(12.5), Inches(7.1), Inches(0.6), Inches(0.3),
             str(page), size=10, color=GREY, align=PP_ALIGN.RIGHT)


def add_table(slide, left, top, width, height, data, *, header=True,
              col_widths=None, font_size=14):
    rows, cols = len(data), len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table
    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = int(width * w / total)
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = ""
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = Inches(0.08)
            tf.margin_top = tf.margin_bottom = Inches(0.05)
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = val
            run.font.size = Pt(font_size)
            run.font.name = "Calibri"
            if r == 0 and header:
                run.font.bold = True
                run.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            else:
                run.font.color.rgb = DARK
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
    return tbl


# 1 — Capa
s = add_slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.0), SW, Inches(0.05))
accent.fill.solid(); accent.fill.fore_color.rgb = WHITE; accent.line.fill.background()
add_text(s, Inches(0.8), Inches(2.3), Inches(11.5), Inches(1.2),
         "RAG sobre Legislação ANEEL", size=54, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(3.5), Inches(11.5), Inches(0.8),
         "Pipeline end-to-end de ingestão, parsing e Q&A jurídico",
         size=24, color=WHITE)
add_text(s, Inches(0.8), Inches(5.4), Inches(11.5), Inches(0.5),
         "Desafio CEIA  ·  27 mil PDFs  ·  6 dias", size=20, color=WHITE)

# 2 — Problema
s = add_slide()
add_header(s, "O Desafio",
           "Construir um \"assistente\" que responde sobre legislação ANEEL")

# Cenário
add_text(s, Inches(0.6), Inches(1.5), Inches(12), Inches(0.4),
         "Imagine a cena", size=22, bold=True, color=NAVY)
add_bullets(s, Inches(0.8), Inches(2.0), Inches(12), Inches(2.0), [
    "27 mil PDFs de leis do setor elétrico (2016 / 2021 / 2022)",
    "Um especialista faz perguntas técnicas — \"qual o prazo da revisão tarifária?\"",
    "Sistema precisa achar a resposta certa, citando a fonte exata",
], size=18)

# 3 caixas — desafios
boxes = [
    ("⚖ Precisão absoluta",
     "Lei errada = problema real. Não pode inventar nem chutar."),
    ("📋 Documentos complexos",
     "Tabelas, artigos revogados, páginas escaneadas — tudo misturado."),
    ("⏱ Pouco tempo, pouco dinheiro",
     "6 dias. Solução cara (OCR + LLM em tudo) custaria centenas de dólares."),
]
for i, (title, desc) in enumerate(boxes):
    x = Inches(0.6 + i * 4.2)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             x, Inches(4.3), Inches(4.0), Inches(2.0))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = NAVY
    add_text(s, x + Inches(0.15), Inches(4.45), Inches(3.7), Inches(0.5),
             title, size=18, bold=True, color=NAVY)
    add_text(s, x + Inches(0.15), Inches(4.95), Inches(3.7), Inches(1.4),
             desc, size=14, color=DARK)

add_text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
         "→ Solução: pipeline esperto que escolhe a ferramenta certa pra cada página.",
         size=16, bold=True, color=NAVY)
add_footer(s, 2)

# 3 — Arquitetura
s = add_slide()
add_header(s, "A Linha de Montagem",
           "Cada PDF passa por 6 etapas até virar resposta")

steps = [
    ("1. Baixar", "Pegar todos os PDFs do site da ANEEL"),
    ("2. Classificar", "Ler é fácil ou vai precisar \"olho humano\"?"),
    ("3. Extrair texto", "Tirar o conteúdo do PDF → texto puro + tabelas"),
    ("4. Cortar em pedaços", "Quebrar texto longo em \"trechos\" de tamanho útil"),
    ("5. Indexar", "Transformar texto em números → guardar em \"banco de busca\""),
    ("6. Responder", "Buscar trechos relevantes → LLM monta a resposta"),
]
top = Inches(1.5)
h = Inches(0.75)
gap = Inches(0.12)
for i, (label, desc) in enumerate(steps):
    y = top + i * (h + gap)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(0.6), y, Inches(3.0), h)
    box.fill.solid(); box.fill.fore_color.rgb = NAVY
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = label
    r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = WHITE
    add_text(s, Inches(3.85), y + Inches(0.18), Inches(9.0), h,
             desc, size=16, color=DARK)
add_footer(s, 3)

# 4 — Ingestão & Parsing
s = add_slide()
add_header(s, "Baixar e \"Ler\" os PDFs",
           "Como tirar o texto de 27 mil documentos sem gastar fortuna")

# Lado esquerdo — download
add_text(s, Inches(0.6), Inches(1.35), Inches(6), Inches(0.4),
         "🌐 Download — Cloudflare bloqueia bots", size=18, bold=True, color=NAVY)
add_bullets(s, Inches(0.8), Inches(1.85), Inches(5.8), Inches(2.0), [
    "Site da ANEEL inspeciona JA3 (TLS fingerprint)",
    "requests / httpx → 403. curl_cffi com impersonate=chrome resolve",
    "Async com 16 conexões paralelas + retry exponencial",
], size=12)

# Lado esquerdo — router por página
add_text(s, Inches(0.6), Inches(3.7), Inches(6), Inches(0.4),
         "🧠 Router por página (não por doc)", size=18, bold=True, color=NAVY)
add_bullets(s, Inches(0.8), Inches(4.2), Inches(5.8), Inches(2.8), [
    "Cada página classificada via PyMuPDF antes do parse:",
    (1, "digital → texto extraível direto"),
    (1, "scanned → poucos chars + alta razão de imagem → vision"),
    (1, "complex → tabelas detectadas → Camelot ou vision"),
    "Decisão por página: doc misto é a regra, não exceção",
], size=12)

# Lado direito — 2 caminhos com ferramentas nomeadas
add_text(s, Inches(7.0), Inches(1.35), Inches(6), Inches(0.4),
         "Fast vs Slow path", size=18, bold=True, color=NAVY)

fast = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(7.0), Inches(1.85), Inches(5.8), Inches(2.2))
fast.fill.solid(); fast.fill.fore_color.rgb = RGBColor(0xE6, 0xF4, 0xE6)
fast.line.color.rgb = RGBColor(0x22, 0x88, 0x44)
add_text(s, Inches(7.2), Inches(1.95), Inches(5.5), Inches(0.5),
         "✓ Fast path (~95% das páginas)",
         size=15, bold=True, color=RGBColor(0x22, 0x88, 0x44))
add_text(s, Inches(7.2), Inches(2.45), Inches(5.5), Inches(1.7),
         "PyMuPDF → texto + bbox dos spans\n"
         "pdfplumber → tabelas simples (markdown)\n"
         "Camelot → tabelas com bordas (lattice)\n"
         "Custo: 0 · totalmente local",
         size=12, color=DARK)

slow = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(7.0), Inches(4.20), Inches(5.8), Inches(2.2))
slow.fill.solid(); slow.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xE0)
slow.line.color.rgb = RGBColor(0xC2, 0x6B, 0x00)
add_text(s, Inches(7.2), Inches(4.30), Inches(5.5), Inches(0.5),
         "⚠ Slow path (~5% das páginas)",
         size=15, bold=True, color=RGBColor(0xC2, 0x6B, 0x00))
add_text(s, Inches(7.2), Inches(4.80), Inches(5.5), Inches(1.7),
         "Gemini 2.5 Flash via Vertex AI (multimodal)\n"
         "Render 180 DPI → prompt que exige ~~tachado~~\n"
         "thinking_budget=0 → custo de input ~$0,15/M tokens\n"
         "Fallback: gemini-2.5-flash-lite",
         size=12, color=DARK)

# Bottom — economia
add_text(s, Inches(0.6), Inches(6.65), Inches(12), Inches(0.5),
         "💰 Vision em 100% das páginas: ~$300+ · com router por página: ~$22 · ~93% de economia",
         size=14, bold=True, color=RGBColor(0x22, 0x88, 0x44))
add_footer(s, 4)

# 5 — Texto revogado: o conceito
s = add_slide()
add_header(s, "Texto Revogado (Strikethrough)",
           "O que é e por que importa")

# Box conceito
add_text(s, Inches(0.6), Inches(1.4), Inches(7.5), Inches(0.4),
         "O que é", size=22, bold=True, color=NAVY)
add_text(s, Inches(0.8), Inches(1.85), Inches(7.3), Inches(0.5),
         "Texto cortado por linha horizontal:",
         size=17, color=DARK)
ex = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.8), Inches(2.4), Inches(7.3), Inches(0.9))
ex.fill.solid(); ex.fill.fore_color.rgb = LIGHT
ex.line.color.rgb = NAVY
tf = ex.text_frame; tf.margin_left = Inches(0.15)
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Art. 5º A distribuidora "
r.font.size = Pt(18); r.font.color.rgb = DARK
r2 = p.add_run()
r2.text = "deve apresentar plano em 30 dias"
r2.font.size = Pt(18); r2.font.strike = True
r2.font.color.rgb = RGBColor(0xB0, 0x20, 0x20)

add_text(s, Inches(0.6), Inches(3.6), Inches(7.5), Inches(0.4),
         "Em legislação ANEEL", size=22, bold=True, color=NAVY)
add_bullets(s, Inches(0.8), Inches(4.05), Inches(7.3), Inches(2.5), [
    "Indica artigo revogado por norma posterior",
    "Texto continua impresso (rastreabilidade)",
    "Mas não vale mais como norma vigente",
], size=16)

# Box "por que importa"
warn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(8.4), Inches(1.4), Inches(4.4), Inches(5.4))
warn.fill.solid(); warn.fill.fore_color.rgb = RGBColor(0xFF, 0xF4, 0xE0)
warn.line.color.rgb = RGBColor(0xC2, 0x6B, 0x00)
add_text(s, Inches(8.6), Inches(1.55), Inches(4.0), Inches(0.5),
         "⚠ Por que importa", size=20, bold=True,
         color=RGBColor(0xC2, 0x6B, 0x00))
add_bullets(s, Inches(8.6), Inches(2.1), Inches(4.0), Inches(4.5), [
    "LLM lê tachado como vigente → resposta errada",
    "Especialista do setor detecta erro instantâneo",
    "Legislação tem zero tolerância a erro citacional",
    "Campo \"situacao\" do metadado descreve o doc inteiro, não trechos internos",
], size=14)

# Problema técnico
add_text(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
         "Problema: PDF não armazena flag \"isto é tachado\". "
         "Texto e linha são camadas separadas — preciso cruzar.",
         size=15, bold=True, color=NAVY)
add_footer(s, 5)

# 6 — Texto revogado: detecção
s = add_slide()
add_header(s, "Como Detectamos o Strikethrough",
           "Cruzar geometria de texto × linhas desenhadas")

# Diagrama topo — 3 testes
add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.4),
         "PDF digital — 3 testes geométricos por span de texto",
         size=20, bold=True, color=NAVY)

tests = [
    ("1. É linha mesmo?",
     "horizontal fina (≤1.2px) e comprida (>4px) — descarta ruído"),
    ("2. Cai no MEIO do texto?",
     "centro da linha na banda 45% central → descarta sublinhado/overline"),
    ("3. Cobre o texto?",
     "sobreposição horizontal ≥ 55% da largura → descarta decoração curta"),
]
top = Inches(1.85)
for i, (title, desc) in enumerate(tests):
    y = top + i * Inches(0.95)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(0.6), y, Inches(7.0), Inches(0.85))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = NAVY
    tf = box.text_frame; tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = desc
    r2.font.size = Pt(13); r2.font.color.rgb = DARK

add_text(s, Inches(0.6), Inches(4.85), Inches(7.5), Inches(0.4),
         "3 testes verdadeiros → span marcado revoked=True",
         size=14, bold=True, color=RGBColor(0x22, 0x88, 0x44))

# Visual ASCII-like sobre banda central
diag = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                          Inches(8.0), Inches(1.85), Inches(4.8), Inches(2.8))
diag.fill.solid(); diag.fill.fore_color.rgb = WHITE
diag.line.color.rgb = NAVY
# Texto no diagrama
add_text(s, Inches(8.15), Inches(1.95), Inches(4.5), Inches(0.4),
         "Anatomia do span", size=14, bold=True, color=NAVY)
# topo
add_text(s, Inches(8.15), Inches(2.35), Inches(4.5), Inches(0.3),
         "─── topo (sublinhado: IGNORA)", size=11, color=GREY)
# texto exemplo
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                          Inches(8.2), Inches(2.7), Inches(4.5), Inches(1.0))
band.fill.solid(); band.fill.fore_color.rgb = RGBColor(0xFF, 0xE8, 0xE8)
band.line.fill.background()
add_text(s, Inches(8.25), Inches(2.75), Inches(4.4), Inches(0.4),
         "texto do span", size=14, color=DARK)
# linha vermelha simulando tachado
strike = s.shapes.add_connector(1, Inches(8.3), Inches(3.15),
                                  Inches(12.5), Inches(3.15))
strike.line.color.rgb = RGBColor(0xC0, 0x20, 0x20)
strike.line.width = Pt(2)
add_text(s, Inches(8.25), Inches(3.25), Inches(4.4), Inches(0.4),
         "← banda 45% central = TACHADO", size=11, bold=True,
         color=RGBColor(0xC0, 0x20, 0x20))
add_text(s, Inches(8.15), Inches(3.85), Inches(4.5), Inches(0.3),
         "─── base (overline: IGNORA)", size=11, color=GREY)

# Caminho 2 — annotations
add_text(s, Inches(0.6), Inches(5.35), Inches(12), Inches(0.4),
         "Caminho extra (digital): annotations com subtype=StrikeOut → atalho",
         size=14, color=DARK)

# PDF escaneado
add_text(s, Inches(0.6), Inches(5.85), Inches(12), Inches(0.4),
         "PDF escaneado", size=20, bold=True, color=NAVY)
add_text(s, Inches(0.8), Inches(6.30), Inches(12), Inches(0.4),
         "Sem geometria → Gemini Vision faz OCR + prompt: "
         "\"trecho tachado deve aparecer como ~~texto~~\"",
         size=14, color=DARK)

# Saída unificada
add_text(s, Inches(0.6), Inches(6.75), Inches(12), Inches(0.4),
         "→ Saída unificada: marcador ~~...~~ propaga até o prompt da LLM "
         "(\"nunca trate como vigente\")",
         size=14, bold=True, color=RGBColor(0x22, 0x88, 0x44))
add_footer(s, 6)

# 7 — Chunking
s = add_slide()
add_header(s, "Cortando o Texto em Pedaços",
           "Por que e como dividir 50 páginas de lei em \"trechos\"")

# Esquerda — analogia
add_text(s, Inches(0.6), Inches(1.4), Inches(6.1), Inches(0.4),
         "📖 A analogia do livro", size=20, bold=True, color=NAVY)
add_text(s, Inches(0.8), Inches(1.9), Inches(5.9), Inches(2.0),
         "LLM não cabe um livro inteiro de uma vez.\n"
         "Precisamos cortar em pedaços (\"chunks\").\n\n"
         "Como cortar?",
         size=15, color=DARK)

# Mau caminho
bad = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                         Inches(0.6), Inches(4.0), Inches(6.1), Inches(1.4))
bad.fill.solid(); bad.fill.fore_color.rgb = RGBColor(0xFF, 0xE8, 0xE8)
bad.line.color.rgb = RGBColor(0xC0, 0x20, 0x20)
add_text(s, Inches(0.8), Inches(4.10), Inches(5.7), Inches(0.4),
         "✗ Jeito ingênuo: cortar a cada 500 palavras",
         size=14, bold=True, color=RGBColor(0xC0, 0x20, 0x20))
add_text(s, Inches(0.8), Inches(4.55), Inches(5.7), Inches(0.9),
         "Corta no meio do Art. 5º. Metade fica num pedaço,\n"
         "metade noutro. Resposta sai pela metade também.",
         size=13, color=DARK)

# Bom caminho
good = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.6), Inches(5.55), Inches(6.1), Inches(1.5))
good.fill.solid(); good.fill.fore_color.rgb = RGBColor(0xE6, 0xF4, 0xE6)
good.line.color.rgb = RGBColor(0x22, 0x88, 0x44)
add_text(s, Inches(0.8), Inches(5.65), Inches(5.7), Inches(0.4),
         "✓ Nosso jeito: cortar pela estrutura da lei",
         size=14, bold=True, color=RGBColor(0x22, 0x88, 0x44))
add_text(s, Inches(0.8), Inches(6.10), Inches(5.7), Inches(1.0),
         "Cada Artigo / Parágrafo / Capítulo vira 1 pedaço.\n"
         "Respeitamos como a lei foi escrita.",
         size=13, color=DARK)

# Direita — payload do chunk
add_text(s, Inches(7.0), Inches(1.4), Inches(6), Inches(0.4),
         "Payload propagado por chunk", size=20, bold=True, color=NAVY)
ex = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(7.0), Inches(1.95), Inches(5.9), Inches(3.6))
ex.fill.solid(); ex.fill.fore_color.rgb = LIGHT
ex.line.color.rgb = NAVY
add_text(s, Inches(7.2), Inches(2.05), Inches(5.5), Inches(0.4),
         "Metadados indexados no Qdrant:",
         size=13, bold=True, color=NAVY)
add_bullets(s, Inches(7.2), Inches(2.50), Inches(5.6), Inches(3.0), [
    "article_ref (Art. 12)",
    "doc_id, arquivo, ano, autor, situacao_doc",
    "page_start, page_end",
    "has_table, has_revoked (booleanos)",
    "ementa, assunto, publicacao",
], size=12)

add_text(s, Inches(7.0), Inches(5.7), Inches(6), Inches(1.4),
         "Permite filtros estruturados na query\n"
         "(ex: ano=2021 AND has_revoked=false) +\n"
         "citação automática [doc_id | art. X].",
         size=12, color=DARK)
add_footer(s, 7)

# 8 — Chunking: estratégia em 3 níveis
s = add_slide()
add_header(s, "Estratégia em 3 Níveis (Cascata)",
           "Não é só regex, nem só recursivo — é os dois trabalhando juntos")

# Topo: introdução
add_text(s, Inches(0.6), Inches(1.35), Inches(12), Inches(0.5),
         "Cada nível só entra em ação se o anterior não resolveu sozinho:",
         size=15, color=DARK)

# Nível 1
n1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.4))
n1.fill.solid(); n1.fill.fore_color.rgb = RGBColor(0xE6, 0xF0, 0xFF)
n1.line.color.rgb = NAVY
# Badge nível
b1 = s.shapes.add_shape(MSO_SHAPE.OVAL,
                        Inches(0.65), Inches(2.15), Inches(0.6), Inches(0.6))
b1.fill.solid(); b1.fill.fore_color.rgb = NAVY
b1.line.fill.background()
tf = b1.text_frame; tf.margin_top = Inches(0.02)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "1"
r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = WHITE

add_text(s, Inches(1.4), Inches(2.10), Inches(11.0), Inches(0.4),
         "Split estrutural por hierarquia legal (regex)",
         size=17, bold=True, color=NAVY)
add_text(s, Inches(1.4), Inches(2.55), Inches(11.0), Inches(0.85),
         "Cabeçalhos: Art. X · § N · Capítulo (romano) · Seção · Título.\n"
         "1 artigo = 1 chunk · article_ref propaga como metadado pra citação.",
         size=13, color=DARK)

# Nível 2
n2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.5), Inches(3.55), Inches(12.3), Inches(1.55))
n2.fill.solid(); n2.fill.fore_color.rgb = RGBColor(0xFF, 0xF4, 0xE0)
n2.line.color.rgb = RGBColor(0xC2, 0x6B, 0x00)
b2 = s.shapes.add_shape(MSO_SHAPE.OVAL,
                        Inches(0.65), Inches(3.70), Inches(0.6), Inches(0.6))
b2.fill.solid(); b2.fill.fore_color.rgb = RGBColor(0xC2, 0x6B, 0x00)
b2.line.fill.background()
tf = b2.text_frame; tf.margin_top = Inches(0.02)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "2"
r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = WHITE

add_text(s, Inches(1.4), Inches(3.65), Inches(11.0), Inches(0.4),
         "SE bloco > 900 tokens → sliding window com overlap 100",
         size=17, bold=True, color=RGBColor(0xC2, 0x6B, 0x00))
add_text(s, Inches(1.4), Inches(4.10), Inches(11.0), Inches(1.0),
         "Limite de 900 escolhido pra caber confortável no max_length=1024 do BGE-M3.\n"
         "Quebra por parágrafo (\\n\\n) preservando coerência · overlap 100 tokens evita\n"
         "perda de contexto na fronteira entre chunks vizinhos.",
         size=13, color=DARK)

# Nível 3
n3 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.5), Inches(5.25), Inches(12.3), Inches(1.4))
n3.fill.solid(); n3.fill.fore_color.rgb = RGBColor(0xE6, 0xF4, 0xE6)
n3.line.color.rgb = RGBColor(0x22, 0x88, 0x44)
b3 = s.shapes.add_shape(MSO_SHAPE.OVAL,
                        Inches(0.65), Inches(5.40), Inches(0.6), Inches(0.6))
b3.fill.solid(); b3.fill.fore_color.rgb = RGBColor(0x22, 0x88, 0x44)
b3.line.fill.background()
tf = b3.text_frame; tf.margin_top = Inches(0.02)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "3"
r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = WHITE

add_text(s, Inches(1.4), Inches(5.35), Inches(11.0), Inches(0.4),
         "SE chunk < 60 tokens → merge no anterior (exceto se has_table=true)",
         size=17, bold=True, color=RGBColor(0x22, 0x88, 0x44))
add_text(s, Inches(1.4), Inches(5.80), Inches(11.0), Inches(0.85),
         "Migalhas (\"§ 3º Revogado.\") inflam o índice e atrapalham reranker.\n"
         "Tabelas mantêm chunk próprio · merge atualiza page_end e has_revoked.",
         size=13, color=DARK)

# Bottom: comparativo curto
add_text(s, Inches(0.6), Inches(6.85), Inches(12), Inches(0.4),
         "📌 Estrutura primeiro · tamanho depois — respeita como a lei foi escrita",
         size=15, bold=True, color=NAVY)
add_footer(s, 8)

# 9 — Embedding & Vector DB
s = add_slide()
add_header(s, "Transformando Texto em \"Coordenadas\"",
           "Como o computador entende o significado de um trecho")

# Topo — modelo escolhido
add_text(s, Inches(0.6), Inches(1.35), Inches(12), Inches(0.4),
         "🧬 Modelo de embedding: BGE-M3",
         size=20, bold=True, color=NAVY)
add_text(s, Inches(0.8), Inches(1.85), Inches(12), Inches(1.0),
         "Cada trecho vira vetor de 1024 dimensões + vetor esparso lexical.\n"
         "Um único modelo entrega os dois — sem rodar BM25 separado.",
         size=14, color=DARK)

# Duas caixas — dense + sparse
b1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.5), Inches(3.0), Inches(6.2), Inches(2.5))
b1.fill.solid(); b1.fill.fore_color.rgb = LIGHT
b1.line.color.rgb = NAVY
add_text(s, Inches(0.7), Inches(3.10), Inches(5.8), Inches(0.5),
         "Vetor denso (1024d, cosine)", size=16, bold=True, color=NAVY)
add_text(s, Inches(0.7), Inches(3.60), Inches(5.8), Inches(1.9),
         "Captura significado.\n"
         "\"prazo de revisão\" ≈ \"período de reanálise\".\n\n"
         "Distância: cosseno · normalizado.",
         size=13, color=DARK)

b2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(6.9), Inches(3.0), Inches(6.0), Inches(2.5))
b2.fill.solid(); b2.fill.fore_color.rgb = LIGHT
b2.line.color.rgb = NAVY
add_text(s, Inches(7.1), Inches(3.10), Inches(5.6), Inches(0.5),
         "Vetor esparso (lexical)", size=16, bold=True, color=NAVY)
add_text(s, Inches(7.1), Inches(3.60), Inches(5.6), Inches(1.9),
         "Captura termos exatos.\n"
         "\"REN 1.000\", \"bandeira tarifária\", siglas.\n\n"
         "Substitui BM25 — nativo no Qdrant.",
         size=13, color=DARK)

# Decisões: por que BGE-M3 e por que Qdrant
add_text(s, Inches(0.5), Inches(5.6), Inches(12), Inches(0.4),
         "Decisões:",
         size=15, bold=True, color=NAVY)
add_bullets(s, Inches(0.7), Inches(6.0), Inches(12.3), Inches(1.4), [
    "BGE-M3 escolhido por: multilíngue forte em PT-BR · 8K tokens de contexto · dense+sparse no mesmo modelo · roda local (custo zero)",
    "Qdrant escolhido por: named vectors nativos (dense + sparse na mesma coleção) · filtros estruturados (ano, doc_id, has_table) · open-source · Docker local",
], size=12)
add_footer(s, 9)

# 10 — Retriever
s = add_slide()
add_header(s, "Retriever Híbrido com RRF + Rerank",
           "3 fontes de busca → fusão → reordenação → top 6")

# Diagrama — 3 fontes → fusão → seleção final
def pill(slide, x, y, w, h, label, color=NAVY, sub=None):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf = box.text_frame; tf.margin_left = Inches(0.05); tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(10); r2.font.color.rgb = WHITE

# Pergunta
pill(s, Inches(0.5), Inches(2.8), Inches(2.0), Inches(0.7),
     "Pergunta", color=DARK)

# 3 fontes
pill(s, Inches(3.1), Inches(1.55), Inches(3.0), Inches(0.85),
     "Dense (BGE-M3)", color=NAVY, sub="top 40 por cosseno")
pill(s, Inches(3.1), Inches(2.7), Inches(3.0), Inches(0.85),
     "Sparse (BGE-M3)", color=NAVY, sub="top 40 lexical")
pill(s, Inches(3.1), Inches(3.85), Inches(3.0), Inches(0.85),
     "Identifier lookup", color=NAVY, sub="regex nº+ano em arquivo")

# Fusão RRF
pill(s, Inches(6.7), Inches(2.7), Inches(2.4), Inches(0.85),
     "RRF (k=60)", color=RGBColor(0x44, 0x66, 0x99),
     sub="fusão de rankings")

# Reranker
pill(s, Inches(9.6), Inches(2.7), Inches(2.2), Inches(0.85),
     "Reranker", color=RGBColor(0x44, 0x66, 0x99),
     sub="Cohere v3 / BGE-v2-m3")

# Top
pill(s, Inches(12.2), Inches(2.7), Inches(0.7), Inches(0.85),
     "Top 6", color=RGBColor(0x33, 0x88, 0x55))

# Decisões
add_text(s, Inches(0.5), Inches(5.05), Inches(12), Inches(0.4),
         "Decisões-chave:", size=16, bold=True, color=NAVY)
add_bullets(s, Inches(0.7), Inches(5.5), Inches(12.3), Inches(2.5), [
    "RRF (Reciprocal Rank Fusion) em vez de pesos fixos: combina rankings sem tuning manual; robusto a escalas diferentes (cosseno vs lexical)",
    "Identifier lookup como 3ª fonte: regex extrai \"REN 1.000/2021\" da query e busca direto no campo arquivo — pergunta por nº de norma não pode depender só de embedding",
    "Hard filter: se identifier match ≤ 3 docs, descarta dense/sparse e busca só dentro deles (evita ruído)",
    "Reranker cross-encoder reordena top 20 → top 6: cohere-rerank-v3 (API) ou BGE-reranker-v2-m3 (local) como fallback",
    "Penalty 0.85 em chunks de ementa (page_start=0): metadados não são conteúdo legal — não devem dominar topo",
], size=11)
add_footer(s, 10)

# 11 — Geração & Avaliação
s = add_slide()
add_header(s, "Resposta e Como Sabemos Se Está Boa",
           "LLM monta a resposta · gabarito mede qualidade")

# Esquerda — geração
add_text(s, Inches(0.6), Inches(1.4), Inches(6.1), Inches(0.4),
         "💬 Como a resposta é montada", size=20, bold=True, color=NAVY)

flow = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.6), Inches(1.95), Inches(6.1), Inches(2.5))
flow.fill.solid(); flow.fill.fore_color.rgb = LIGHT
flow.line.color.rgb = NAVY
add_text(s, Inches(0.8), Inches(2.05), Inches(5.7), Inches(2.4),
         "Modelo: Gemini 2.5 Flash (Vertex AI)\n"
         "1. Top-6 chunks viram contexto\n"
         "2. System prompt impõe:\n"
         "    – citação [doc_id | art. X] obrigatória\n"
         "    – tratar ~~...~~ como revogado\n"
         "    – \"não consta\" quando ausente\n"
         "3. Saída: resposta + payload de citações",
         size=13, color=DARK)

add_text(s, Inches(0.6), Inches(4.65), Inches(6.1), Inches(0.4),
         "Decisões da geração", size=18, bold=True, color=NAVY)
add_bullets(s, Inches(0.7), Inches(5.10), Inches(6.0), Inches(2.0), [
    "temperature=0 + thinking_budget=0 → resposta determinística e barata",
    "Gemini 2.5 Flash sobre Pro: latência 3-4× menor, qualidade suficiente",
    "Fallback automático para gemini-2.5-flash-lite em erro",
], size=11)

# Direita — avaliação
add_text(s, Inches(7.0), Inches(1.4), Inches(6), Inches(0.4),
         "Avaliação — benchmark anotado", size=20, bold=True, color=NAVY)
add_text(s, Inches(7.2), Inches(1.85), Inches(5.9), Inches(0.7),
         "Gold JSONL com expected_doc_ids, expected_articles,\n"
         "numeric_facts, must_contain, should_refuse.",
         size=12, color=DARK)

add_bullets(s, Inches(7.2), Inches(2.85), Inches(5.9), Inches(4.0), [
    "Retrieval: Hit@k · Recall@k · Precision@k · nDCG@k · MRR",
    "Citation-F1 (regex normaliza nº+ano)",
    "Numeric exactness (tolerância em valores)",
    "Refusal correctness (detecta \"não consta\")",
    "LLM-as-judge: faithfulness · correctness · completeness · critical_error",
    "Latência p50/p95 separadas (retrieval vs geração)",
    "Bootstrap CI 95% — comparação estatística entre versões do retriever",
], size=11)
add_footer(s, 11)

# 12 — Resultado & Trade-offs
s = add_slide()
add_header(s, "O Que Funcionou e o Que Foi Trocado",
           "Decisões importantes em linguagem simples")

add_text(s, Inches(0.6), Inches(1.4), Inches(6), Inches(0.4),
         "✓ Acertos", size=22, bold=True,
         color=RGBColor(0x22, 0x88, 0x44))
add_bullets(s, Inches(0.8), Inches(1.95), Inches(5.9), Inches(3.5), [
    "Decisor por página economizou ~93% do custo de OCR",
    "Cortar texto pelos artigos da lei melhorou muito as citações",
    "Pipeline retoma do ponto onde parou se cair — sem refazer do zero",
], size=15)

add_text(s, Inches(7.0), Inches(1.4), Inches(6), Inches(0.4),
         "⚖ Escolhas conscientes", size=22, bold=True,
         color=RGBColor(0xC2, 0x6B, 0x00))
add_bullets(s, Inches(7.2), Inches(1.95), Inches(5.9), Inches(3.5), [
    "Pagamos a LLM (Vertex AI) em vez de usar gratuita — confiabilidade vale mais",
    "Reranker pago é mais rápido e preciso, mas tem limite de chamadas",
    "Banco de controle simples (sqlite) — basta pra esse projeto",
], size=15)

add_text(s, Inches(0.6), Inches(5.5), Inches(12), Inches(0.4),
         "🚀 Próximos passos", size=22, bold=True, color=NAVY)
add_bullets(s, Inches(0.8), Inches(6.05), Inches(12), Inches(1.5), [
    "Testar variações de tamanho de chunk · LLM reescrevendo a pergunta antes de buscar · cache de buscas frequentes",
], size=14)
add_footer(s, 12)

# 13 — Stack final + obrigado
s = add_slide()
add_header(s, "Ferramentas Usadas",
           "Cada etapa, qual ferramenta e pra que serve")
data = [
    ["Etapa", "Ferramenta", "Pra que serve"],
    ["Baixar PDFs", "curl_cffi", "Disfarça de navegador pra passar pelo Cloudflare"],
    ["Ler PDF normal", "PyMuPDF + pdfplumber + Camelot", "Tira texto e tabelas direto do PDF"],
    ["Ler PDF escaneado", "Gemini 2.5 Flash", "LLM \"olha\" a página e transcreve"],
    ["Cortar em pedaços", "Lógica própria", "Quebra pelos artigos da lei"],
    ["DNA do texto", "BGE-M3", "Vira números pra busca por significado"],
    ["Banco de busca", "Qdrant", "Guarda os números e busca rápido"],
    ["Encontrar trechos", "3 buscas + fusão + juiz", "Acha os 6 melhores"],
    ["Gerar resposta", "Gemini 2.5 Flash", "Lê os trechos e responde com citação"],
    ["Avaliar", "Gabarito + outra LLM como juiz", "Mede se a resposta está boa"],
]
add_table(s, Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.4),
          data, col_widths=[3, 4, 6], font_size=12)
add_text(s, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.5),
         "Obrigado!", size=28, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER)
add_footer(s, 13)

out = "c:/Users/Gustavo/Desktop/CEIA/desafio-rag/apresentacao.pptx"
prs.save(out)
print(f"saved: {out}")
